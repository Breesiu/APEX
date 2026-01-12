import os
from ..schema import AgentState
from ..tools.pptx_parser import get_element_shape_map
from pptx import Presentation
from ..tools.pptx_execuator import API_executor
from typing import List
import json
from ..tools import pptx_execuator
from ..tools.image_tools import convert_pptx_to_png
from pptx.util import Inches, Cm, Pt


def code_generation_execution_agent_api(state: AgentState) -> dict:
    """
    Code Generator & Executor Agent: Applies edit operations to PPTX.
    
    This agent doesn't use LLM - it directly executes structured operations.
    """
    logger = state.logger
    logger.info("\n" + "="*60)
    logger.info("AGENT: Code Generation & Execution Agent")
    logger.info("="*60)
    
    # if not state.action_plan:
    #     return {"error": "No action plan available for execution."}
    
    # Determine input and output paths
    input_pptx = state.current_pptx_path or state.pptx_path
    
    iteration = state.iteration_count
    if state.iterate_version and iteration>0:
        # output_filename = f"poster_debug_v{iteration + 1}_{state.iterate_version}_{state.timestamp}.pptx"
        output_filename = f"poster_debug_v{iteration + 1}.pptx"

    else:
        output_filename = f"poster_v{iteration + 1}.pptx"
    # output_path = os.path.join(state.output_dir, output_filename)
    output_path = state.output_dir / output_filename
    # Get element-to-shape mapping
    # element_shape_map = get_element_shape_map()
    
    logger.info(f"\nInput: {input_pptx}")
    logger.info(f"Output: {output_path}")
    # print(f"Operations to apply: {len(state.plan.operations)}")
    
    # excute operations in action plan
    prs = pptx_execuator.get_current_state().prs #Presentation(input_pptx)
    
    pptx_execuator.get_current_state().pptx_folder_path = os.path.dirname(str(state.pptx_path))
    
    slide = prs.slides[0]
    
    # Create reverse map: shape_id -> shape
    shape_map = {s.shape_id: s for s in slide.shapes}
    context = {
        "slide":slide,
        "shape_map": shape_map,
    }
    # Sort operations by priority (higher priority first)
    #sorted_ops = sorted(state.action_plan.operations, key=lambda x: x.priority, reverse=True)
    
    # 构造 API 调用行
    if False:
        api_lines = build_api_lines(state.action_plan.operations)
    else:
        api_lines = state.api_list

    # 执行
    error_info = API_executor(api_lines, api_context=None, prs = prs, logger=logger)
    '''
    for op in sorted_ops:
        try:
            # excute each operation in python, EditOperation includes function, parameters, etc.
            #TODO: one line to apply operations
            
            
            pass
            
        except Exception as e:
            print(f"Warning: Failed to apply operation {op.op_type} on {op.element_id}: {e}")
            return {"error": f"Failed to apply operation {op.op_type} on {op.params}: {str(e)}"} # op.params?
        
    '''
    prs.save(output_path)
    convert_pptx_to_png(output_path)

    
    # print(f"\n✓ Successfully applied {len(state.action_plan.operations)} operations")
    return {
        "current_pptx_path": output_path,
        # "operations_applied": state.action_plan.operations,
        "iteration_count": iteration + 1
    }
    
    
def get_all_pptx_api_functions():
    """
    Automatically collects all public functions from pptx_api.py.
    Only public callables are injected into the LLM execution environment.
    """
    api_dict = {}

    for name in dir(pptx_execuator):
        if name.startswith("_"):
            continue

        obj = getattr(pptx_execuator, name)
        # 仅注入函数
        if callable(obj):
            api_dict[name] = obj

    return api_dict