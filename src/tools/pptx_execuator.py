from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import re
from typing import List, Dict, Optional, Tuple
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

import io
import os
import threading
from typing import Optional
import contextvars
# ============================================================================
# 全局状态管理
# ============================================================================

class PosterState:
    """管理海报的全局状态"""

    def __init__(self):
        self.prs = None
        self.slide = None
        self.shape_map = {}  # element_(str) -> shape
        self.current_shape = None
        self._next_element_id = 1  # 自增元素ID计数器
        self.pptx_folder_path = ""

    def set_from_prs(self, prs):
        """从 Presentation 对象设置状态"""
        self.prs = prs
        self.slide = prs.slides[0]  # 单页海报
        self._build_shape_map()
        return f"Loaded poster with {len(self.shape_map)} elements"

    def load_poster(self, ppt_path: str):
        """加载海报文件"""
        self.pptx_folder_path = os.path.dirname(ppt_path)
        self.prs = Presentation(ppt_path)
        return self.set_from_prs(self.prs)

    def _build_shape_map(self):
        """构建元素映射表 (使用 shape.name 作为 Key)"""
        # self.shape_map = _element_to_shape_map
        # # for shape in self.slide.shapes:
        # #     # 使用 name 作为键，如果 name 为空，通常 PPT 会自动分配，但为了保险可以做个处理
        # #     key = shape.name if shape.name else str(shape.shape_id)
        # #     # key = shape._element.get('mytag') if shape._element.get('mytag') else str(shape.shape_id)
        # #     self.shape_map[key] = shape
        # #     print(key)
        
        
        # Already make it in parse_pptx_to_json in pptx_parser.py
        # for key, value in self.shape_map.items():
        #     print(key)
            
        
        if self.shape_map:
            self._next_element_id = max(int(k) for k in self.shape_map.keys()) + 1

        else:
            self._next_element_id = 1

    def save_poster(self, output_path: str):
        """保存海报"""
        self.prs.save(output_path)
        return f"Poster saved to {output_path}"


    def get_shape(self, element_id: str):
        """根据 ID 获取 shape"""
        if element_id not in self.shape_map:
            raise ValueError(f"Element '{element_id}' not found. Available: {list(self.shape_map.keys())}")
        return self.shape_map[element_id]

    # def get_shape(self, element_id: str):
    #     for shape in self.shape_map.values():
    #         if shape.name == element_id:
    #             return shape
    #     raise ValueError(f"Element with name {shape.name} not found.")
    

_state_context_var = contextvars.ContextVar("poster_state", default=None)
def get_current_state() -> PosterState:
    """获取当前协程/任务的状态"""
    state = _state_context_var.get()
    
    # 如果当前上下文中没有状态，则初始化一个
    if state is None:
        state = PosterState()
        _state_context_var.set(state)
        
    return state
def set_current_state(state: PosterState):
    """设置或更新当前上下文的状态"""
    _state_context_var.set(state)
    
    
# ============================================================================
# 颜色工具
# ============================================================================
from pptx.dml.color import RGBColor
from PIL import ImageColor

def hex_to_rgb(color_str: str) -> RGBColor:
    """
    将颜色名称或 Hex 转换为 pptx RGBColor
    支持输入: "red", "Light Gray", "Deep Sky Blue", "#FF0000", "rgb(255, 0, 0)"
    """
    # 1. 基础清理：去除首尾空格，转小写（虽然 ImageColor 不区分大小写，但方便处理）
    raw_color = color_str.strip()
    
    try:
        # 尝试直接转换
        rgb_tuple = ImageColor.getrgb(raw_color)
    except ValueError:
        try:
            # 2. 容错处理：如果是 "light gray" 这种带空格的名称
            # 只有当它不包含 '(' (即不是 rgb(...) 格式) 时才去空格
            # 这样可以把 "Light Gray" 变成 "lightgray"，但保留 "rgb(255, 0, 0)" 的结构
            if "(" not in raw_color:
                normalized_color = raw_color.replace(" ", "")
                rgb_tuple = ImageColor.getrgb(normalized_color)
            else:
                raise ValueError # 如果是 rgb() 格式出错，直接抛出
        except ValueError:
            # 3. 最终失败回退
            print(f"Warning: Unknown color '{color_str}', fallback to black.")
            return RGBColor(0, 0, 0)

    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])



def _apply_markdown_bold(paragraph, text):
    parts = [p for p in re.split(r'(\*\*.*?\*\*)', text) if p]

    base_font = None
    if paragraph.runs:
        base_font = paragraph.runs[-1].font

    for part in parts:
        run = paragraph.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part
            if base_font:
                run.font.bold = base_font.bold

        if base_font:
            run.font.name = base_font.name
            run.font.size = base_font.size
            try:
                if base_font.color.rgb:
                    run.font.color.rgb = base_font.color.rgb
            except Exception:
                pass


# ============================================================================
# 1. 文本格式 API - 字体大小、颜色、加粗
# ============================================================================

def set_text_font_size(element_id: str, font_size: float):
    """
    设置文本字体大小

    Args:
        element_id: 元素ID
        size: 字号（磅）
        paragraph_idx: 段落索引，None表示所有段落
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

    return f"Set font size to {font_size}pt for element {element_id}"


def set_text_color(element_id: str, color: str):
    """
    设置文本颜色

    Args:
        element_id: 元素ID
        color: 颜色名称或十六进制值
        paragraph_idx: 段落索引，None表示所有段落
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    rgb_color = hex_to_rgb(color)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = rgb_color

    return f"Set text color to {color} for element {element_id}"


def set_text_bold(element_id: str, bold: bool = True):
    """
    设置文本加粗

    Args:
        element_id: 元素ID
        bold: True为加粗，False为取消加粗
        paragraph_idx: 段落索引，None表示所有段落
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = bold

    return f"Set bold={bold} for element {element_id}"


def set_text_italic(element_id: str, italic: bool = True):
    """设置文本斜体"""
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.italic = italic

    return f"Set italic={italic} for element {element_id}"


def set_text_underline(element_id: str, underline: bool = True):
    """设置文本下划线"""
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.underline = underline

    return f"Set underline={underline} for element {element_id}"


def highlight_keywords(element_id: str, keywords: List[str], color: Optional[str] = "black",
                       bold: Optional[bool] = True):
    """
    只加粗/变色关键词，而不会影响 run 中的其他文本
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    rgb = hex_to_rgb(color)
    text_frame = shape.text_frame
    count = 0

    for paragraph in text_frame.paragraphs:
        new_runs = []

        for run in paragraph.runs:
            text = run.text
            lower_text = text.lower()

            # 记录当前 run 是否包含关键词
            split_positions = []

            for kw in keywords:
                kw_lower = kw.lower()
                start = 0
                while True:
                    idx = lower_text.find(kw_lower, start)
                    if idx == -1:
                        break
                    split_positions.append((idx, idx + len(kw), kw))
                    start = idx + len(kw)

            # 如果 run 中没有关键词 → 直接加入
            if not split_positions:
                new_runs.append(("normal", text, run))
                continue

            # 如果有多个 keyword 出现，要按起始位置排序
            split_positions.sort(key=lambda x: x[0])

            cursor = 0
            for start, end, kw in split_positions:
                # 1. 关键词前的文本
                if start > cursor:
                    new_runs.append(("normal", text[cursor:start], run))

                # 2. 关键词文本
                new_runs.append(("highlight", text[start:end], run))
                count += 1

                cursor = end

            # 3. 关键词后的文本
            if cursor < len(text):
                new_runs.append(("normal", text[cursor:], run))

        # 🔥 清空原 runs，重新添加拆分后的 runs
        paragraph.clear()

        for run_type, txt, original_run in new_runs:
            new_run = paragraph.add_run()
            new_run.text = txt

            # 复制原 run 的格式（不包括颜色/加粗）
            new_run.font.size = original_run.font.size
            new_run.font.name = original_run.font.name

            if run_type == "highlight":
                new_run.font.bold = bold
                new_run.font.color.rgb = rgb

    return f"Highlighted {count} occurrences in element {element_id}"


def set_font_name(element_id: str, font_name: str):
    """
    设置字体族

    Args:
        element_id: 元素ID
        font_name: 字体名称（如 'Arial', 'Times New Roman'）
        paragraph_idx: 段落索引
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name

    return f"Set font family to {font_name} for element {element_id}"


def _format_substrings(element_id: str, keywords: List[str], 
                       font_size: Optional[float] = None, 
                       color: Optional[str] = None,
                       bold: Optional[bool] = None, 
                       italic: Optional[bool] = None, 
                       underline: Optional[bool] = None,
                       font_name: Optional[str] = None):
    
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    rgb = hex_to_rgb(color) if color else None
    text_frame = shape.text_frame
    count = 0

    for paragraph in text_frame.paragraphs:
        new_runs = []
        
        for run in paragraph.runs:
            text = run.text
            if not text:
                continue
            
            lower_text = text.lower()
            matches = []
            
            for kw in keywords:
                if not kw: continue
                kw_lower = kw.lower()
                start = 0
                while True:
                    idx = lower_text.find(kw_lower, start)
                    if idx == -1:
                        break
                    matches.append((idx, idx + len(kw)))
                    start = idx + len(kw)
            
            matches.sort(key=lambda x: x[0])
            
            merged_matches = []
            if matches:
                curr_start, curr_end = matches[0]
                for next_start, next_end in matches[1:]:
                    if next_start < curr_end:
                        curr_end = max(curr_end, next_end)
                    else:
                        merged_matches.append((curr_start, curr_end))
                        curr_start, curr_end = next_start, next_end
                merged_matches.append((curr_start, curr_end))
            
            cursor = 0
            for start, end in merged_matches:
                if start > cursor:
                    new_runs.append(("normal", text[cursor:start], run))
                
                new_runs.append(("highlight", text[start:end], run))
                count += 1
                cursor = end
            
            if cursor < len(text):
                new_runs.append(("normal", text[cursor:], run))

        paragraph.clear()
        
        for run_type, txt, original_run in new_runs:
            new_run = paragraph.add_run()
            new_run.text = txt
            
            if original_run.font.size: new_run.font.size = original_run.font.size
            if original_run.font.name: new_run.font.name = original_run.font.name
            if original_run.font.bold is not None: new_run.font.bold = original_run.font.bold
            if original_run.font.italic is not None: new_run.font.italic = original_run.font.italic
            if original_run.font.underline is not None: new_run.font.underline = original_run.font.underline
            try:
                if original_run.font.color.type: 
                     new_run.font.color.rgb = original_run.font.color.rgb
            except:
                pass

            if run_type == "highlight":
                if font_size: new_run.font.size = Pt(font_size)
                if font_name: new_run.font.name = font_name
                if bold is not None: new_run.font.bold = bold
                if italic is not None: new_run.font.italic = italic
                if underline is not None: new_run.font.underline = underline
                if rgb: new_run.font.color.rgb = rgb

    return f"Formatted {count} occurrences of keywords in element {element_id}"


def text_format_brush(element_id: str, font_size: Optional[float] = None, color: Optional[str] = None,
                      bold: Optional[bool] = None, italic: Optional[bool] = None, underline: Optional[bool] = None,
                      font_name: Optional[str] = None, words: Optional[List[str]] = None):
    
    if words is not None:
        # Split text into words to handle potential spacing/newline mismatches
        # keywords = text.split()
        return _format_substrings(element_id, words, font_size, color, bold, italic, underline, font_name)

    if font_size:
        set_text_font_size(element_id, font_size)
    if color:
        set_text_color(element_id, color)
    if bold:
        set_text_bold(element_id, bold)
    if italic:
        set_text_italic(element_id, italic)
    if underline:
        set_text_underline(element_id, underline)
    if font_name:
        set_font_name(element_id, font_name)
    
    return f"Formatted element {element_id}"


# ============================================================================
# 2. 通用的位置和大小 API
# ============================================================================
def set_element_size(element_id: str, width: Optional[float] = None, height: Optional[float] = None):
    """
    设置元素大小

    Args:
        element_id: 元素ID
        width: 宽度
        height: 高度
    """
    shape = get_current_state().get_shape(element_id)

    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)

    return f"Set image sizefor element {element_id}"


def set_element_position(element_id: str, left: Optional[float]=None, top: Optional[float]=None):
    """
    设置元素位置

    Args:
        element_id: 元素ID
        left: 左边距
        top: 上边距
    """

    shape = get_current_state().get_shape(element_id)

    if left:
        shape.left = Inches(left)
    if top:
        shape.top = Inches(top)

    return f"Set position for {element_id}"


def resize_element_proportionally(element_id: str, scale: float, fixed_center: Optional[bool] = False):
    """
    按比例缩放

    Args:
        element_id: 元素ID
        scale: 缩放比例（如 1.5 表示放大1.5倍）
        fixed_center: 是否固定中心点缩放。True为固定中心，False为固定左上角。
    """
    shape = get_current_state().get_shape(element_id)

    old_width = shape.width
    old_height = shape.height

    new_width = int(old_width * scale)
    new_height = int(old_height * scale)

    if fixed_center:
        # Calculate center
        center_x = shape.left + old_width / 2
        center_y = shape.top + old_height / 2

        # Update dimensions
        shape.width = new_width
        shape.height = new_height

        # Update position to maintain center
        shape.left = int(center_x - new_width / 2)
        shape.top = int(center_y - new_height / 2)
    else:
        shape.width = new_width
        shape.height = new_height

    return f"Resized element {element_id} by {scale}x "


def move_element_relative(element_id: str, delta_x: Optional[float] = 0, delta_y: Optional[float] = 0):
    """
    相对移动

    Args:
        element_id: 元素ID
        delta_x: X轴移动距离
        delta_y: Y轴移动距离
    """
    shape = get_current_state().get_shape(element_id)

    shape.left += Inches(delta_x)
    shape.top += Inches(delta_y)

    return f"Moved image for element {element_id}"


# ============================================================================
# 3. 形状 API
# ============================================================================

def _get_dash_style(style_name: str):
    if not style_name:
        return None
    dash_map = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dash": MSO_LINE_DASH_STYLE.DASH,  # . . . .
        "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,  # - . - .
        "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,  # —— ——
    }
    return dash_map.get(style_name.lower(), MSO_LINE_DASH_STYLE.SOLID)


def insert_shape(left: float, top: float, width: float, height: float, shape_type: str,
                 fill_color: Optional[str] = None, line_color: Optional[str] = None,
                 line_width: float = 0.0, line_dash: Optional[str] = None, element_id: Optional[str] = None):
    """

    Args:
        left, top, width, height: 位置和大小（厘米）
        fill_color: 填充颜色（HEX）
    """

    type_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "diamond": MSO_SHAPE.DIAMOND,
        "oval": MSO_SHAPE.OVAL,
        "star": MSO_SHAPE.STAR_5_POINT,
        "curved right arrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
        "curved left arrow": MSO_SHAPE.CURVED_LEFT_ARROW
    }

    # 创建 shape
    shape = get_current_state().slide.shapes.add_shape(
        type_map.get(shape_type.lower()),
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    # 设置填充颜色
    if fill_color:
        rgb_color = hex_to_rgb(fill_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_color
    else:
        shape.fill.background()  # 设置为无填充
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = hex_to_rgb(line_color)
        if line_width > 0:
            shape.line.width = Pt(line_width)
        if line_dash:
            shape.line.dash_style = _get_dash_style(line_dash)
    else:
        shape.line.fill.background()  # 设置为无轮廓
    
    if not element_id:
        element_id = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1
    get_current_state().shape_map[element_id] = shape
    shape.name = element_id

    return f"{element_id}"


def insert_line(start_x: float, start_y: float, end_x: float, end_y: float,
                color: Optional[str] = "black", width: Optional[float] = 1.0, dash_style: Optional[str] = 'solid',
                element_id: Optional[str] = None):
    """
    插入一条直线

    Args:
        start_x, start_y: 起点坐标
        end_x, end_y: 终点坐标
        color: 线条颜色
        width: 线条宽度 (磅)

    Returns:
        新线条的 element_id
    """

    # 创建线条
    shape = get_current_state().slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )

    # 设置线条颜色
    rgb_color = hex_to_rgb(color)
    line = shape.line
    line.color.rgb = rgb_color
    line.width = Pt(width)

    if dash_style and dash_style != "solid":
        shape.line.dash_style = _get_dash_style(dash_style)

    if not element_id:
        element_id = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1
    get_current_state().shape_map[element_id] = shape
    shape.name = element_id

    return f"{element_id}"

from pptx.oxml.ns import qn

def change_rect_to_rounded_rect(shape):
    # 1. 获取形状底层的 XML 元素 (spPr -> prstGeom)
    sp_pr = shape.element.spPr
    prst_geom = sp_pr.prstGeom
    
    # 2. 检查当前是否为预设几何形状
    if prst_geom is None:
        print("该形状不是预设几何形状，无法转换。")
        return

    # 3. 关键一步：直接修改 XML 属性 'prst'
    current_type = prst_geom.get('prst')
    if current_type != 'roundRect':
        prst_geom.set('prst', 'roundRect')
    

# 使用示例
# change_rect_to_rounded_rect(my_shape_object)

def set_shape_style(element_id: str, fill_color: Optional[str] = None, shape_type: Optional[str] = None,
                    line_color: Optional[str] = None,
                    line_width: Optional[float] = None,
                    line_dash: Optional[str] = None):
    shape = get_current_state().get_shape(element_id)
    if fill_color:
        rgb_color = hex_to_rgb(fill_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_color
    
    if shape_type == 'rounded_rectangle':
        change_rect_to_rounded_rect(shape)

    if hasattr(shape, 'line'):
        if line_color or line_width or line_dash:
            set_line_style(element_id, line_color, line_width, line_dash)

    return f"Set style for element {element_id}"


def set_line_style(element_id: str, color: Optional[str] = None,
                   width: Optional[float] = None, dash_style: Optional[str] = None):
    """
    修改现有元素（线条或形状边框）的线条样式。
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'line'):
        return f"Error: Element {element_id} does not support line styling"

    # 1. 修改颜色
    if color:
        rgb_color = hex_to_rgb(color)
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = rgb_color
        # shape.line.color.rgb = rgb_color
    # 2. 修改宽度
    if width is not None:
        shape.line.width = Pt(width)

    # 3. 修改虚线样式
    if dash_style:
        style_enum = _get_dash_style(dash_style)
        if style_enum:
            shape.line.dash_style = style_enum

    return f"Updated line style for {element_id}"


# ============================================================================
# 4. 排版 API
# ============================================================================

def align_elements_x_axis(element_ids: List[str], alignment: str = "left", reference_id: Optional[str] = None):
    """
    水平方向对齐元素（控制左右位置）。

    Args:
        element_ids: 需要移动的元素ID列表。
        alignment: 对齐方式，可选 'left', 'center', 'right'。
        reference_id: 基准元素ID。如果提供，所有元素将与该元素对齐；
                      如果不提供，默认以 element_ids 中的第一个元素作为基准。
    """
    if not element_ids:
        return "No elements to align."

    # 1. 确定基准元素 (Reference Shape)
    # 如果指定了 reference_id，用指定的；否则默认用列表第一个
    target_ref_id = reference_id if reference_id else element_ids[0]
    ref_shape = get_current_state().get_shape(target_ref_id)

    # 获取基准线的坐标值
    ref_left = ref_shape.left
    ref_center = ref_shape.left + (ref_shape.width / 2)
    ref_right = ref_shape.left + ref_shape.width

    shapes = [get_current_state().get_shape(eid) for eid in element_ids]

    # 2. 执行对齐逻辑
    for shape in shapes:
        # 如果是基准元素自己，可以选择跳过，或者重算一遍也无妨
        if alignment == "left":
            # 左对齐：所有元素的 left 等于基准的 left
            shape.left = ref_left

        elif alignment == "center":
            # 居中对齐：元素的中心点 等于 基准的中心点
            # 新 left = 基准中心 - (自身宽度 / 2)
            shape.left = int(ref_center - (shape.width / 2))

        elif alignment == "right":
            # 右对齐：元素的右侧 等于 基准的右侧
            # 新 left = 基准右侧 - 自身宽度
            shape.left = int(ref_right - shape.width)

    return f"Aligned elements horizontally"


def align_elements_y_axis(element_ids: List[str], alignment: str = "top", reference_id: Optional[str] = None):
    """
    垂直方向对齐元素（控制上下位置）。

    Args:
        element_ids: 需要移动的元素ID列表。
        alignment: 对齐方式，可选 'top', 'middle', 'bottom'。
        reference_id: 基准元素ID。如果提供，所有元素将与该元素对齐。
    """
    if not element_ids:
        return "No elements to align."

    target_ref_id = reference_id if reference_id else element_ids[0]
    ref_shape = get_current_state().get_shape(target_ref_id)

    ref_top = ref_shape.top
    ref_middle = ref_shape.top + (ref_shape.height / 2)
    ref_bottom = ref_shape.top + ref_shape.height

    shapes = [get_current_state().get_shape(eid) for eid in element_ids]

    for shape in shapes:
        if alignment == "top":
            # 顶端对齐
            shape.top = ref_top

        elif alignment == "middle":
            # 垂直居中
            shape.top = int(ref_middle - (shape.height / 2))

        elif alignment == "bottom":
            # 底端对齐
            shape.top = int(ref_bottom - shape.height)

    return f"Aligned {len(element_ids)} elements vertically "

def set_text_alignment(element_id: str, alignment: str = "left"):
    """
    设置文本对齐方式
    
    Args:
        element_id: 元素ID
        alignment: 'left', 'center', 'right', 'justify'
    """
    shape = get_current_state().get_shape(element_id)
    
    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"
    
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    align_value = alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)
    
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = align_value
    
    return f"Set text alignment to {alignment} for element {element_id}"
# ============================================================================
# 5. 文本内容编辑 API - 增删改
# ============================================================================
'''
def set_text_content(element_id: str, text: str, font_size: Optional[float] = 44, font_name: Optional[str] = "Arial",
                     color: Optional[str] = 'black', bold: Optional[bool] = False, italic: Optional[bool] = False,
                     underline: Optional[bool] = False):
    """
    设置文本内容（完全替换），并处理项目符号。

    Args:
        element_id: 元素ID
        text: 新文本内容
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"
    
    tf = shape.text_frame
    tf.clear()  # 清除所有现有段落和格式

    # 保证至少有一个段落用于后续操作
    if not tf.paragraphs:
        tf.add_paragraph()
    #print(f"Setting text content for element {element_id}: {clean_text}")
    # 删除默认创建的空段落
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text.strip():
        p = tf.paragraphs[0]
        p.clear()
        # We need to remove the paragraph itself if it's empty.
        # A bit of a hack to get the parent and remove the child element.
        if len(p.runs) == 0:
            p_element = p._p
            p_element.getparent().remove(p_element)

    lines = text.split('\n')
    for line in lines:
        p = tf.add_paragraph()
        stripped_line = line.lstrip()  # Use lstrip to preserve indentation on the right

        if stripped_line.startswith('•'):
            p.level = 0
        elif stripped_line.startswith('◦'):
            p.level = 1

        _apply_markdown_bold(p, line)  # Pass the original line to preserve content
    text_format_brush(element_id=element_id, font_size=font_size, font_name=font_name, color=color, bold=bold,
                      italic=italic, underline=underline)
    return f"Set text content for element {element_id}"
'''
def set_text_content(element_id: str, text: str, font_size: Optional[float] = 44, font_name: Optional[str] = "Arial",
                     color: Optional[str] = 'black', bold: Optional[bool] = False, italic: Optional[bool] = False,
                     underline: Optional[bool] = False):
    """
    设置文本内容（完全替换），并根据内容自动调整高度，处理各种换行符。
    """
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    if not isinstance(text, str):
        text = str(text) if text is not None else ""

    # =========================================================
    # Part 1: 文本清洗 (标准化换行符)
    # =========================================================
    # 顺序很关键：先处理字面量 \\n，再处理 \r\n，最后处理 \r
    clean_text = text.replace('\\n', '\n').replace('\r\n', '\n').replace('\r', '\n')
    
    # 这一步很重要：原代码里你 split 的是 text，这里必须 split clean_text
    lines = clean_text.split('\n')

    # =========================================================
    # Part 2: 高度自动估算 (核心新增逻辑)
    # =========================================================
    # 获取当前文本框的宽度 (英寸)
    current_width_inch = shape.width.inches
    
    # 确定计算用的字号 (如果传入 None，默认给一个兜底值，这里取函数默认值 44)
    calc_font_size = font_size if font_size is not None else 44.0
    
    avg_char_width_pt = calc_font_size * 0.55
    
    # 计算有效内容宽度 (减去左右内边距 buffer，约 10pt)
    safe_width_pt = (current_width_inch * 72) - 10
    if safe_width_pt < 10: safe_width_pt = 10
    
    chars_per_line = safe_width_pt / avg_char_width_pt
    if chars_per_line < 1: chars_per_line = 1

    # 计算总折行数
    total_lines = 0
    for line in lines:
        length = len(line)
        if length == 0:
            total_lines += 1 # 空段落算一行
        else:
            # 向上取整：计算该段落会自动折成几行
            total_lines += math.ceil(length / chars_per_line)

    # 计算目标高度 (字号 * 1.2倍行距 * 行数 + 上下缓冲)
    line_height_pt = calc_font_size*1.2
    new_height_inch = ((total_lines * line_height_pt) / 72.0) + 0.2
    
    # 【应用高度】立即更新 Python 对象属性
    shape.height = Inches(new_height_inch)

    # =========================================================
    # Part 3: 写入内容 (保持原有逻辑优化)
    # =========================================================
    tf = shape.text_frame
    tf.clear()  # 清除所有段落
    
    # 设置自动调整属性 (配合 XML 修复使用)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # 删除默认生成的空白段落 (如果存在)
    # 这里的逻辑比原版更简洁一些，直接 check 并 remove
    if len(tf.paragraphs) > 0:
         p = tf.paragraphs[0]
         p._p.getparent().remove(p._p)

    # 写入新内容
    for line in lines:
        p = tf.add_paragraph()
        stripped_line = line.lstrip()

        # 处理 Markdown 列表符号
        if stripped_line.startswith('•') or stripped_line.startswith('- '):
            p.level = 0
        elif stripped_line.startswith('◦') or stripped_line.startswith('  -'):
            p.level = 1

        _apply_markdown_bold(p, line)

    # =========================================================
    # Part 4: 格式应用与 XML 修复
    # =========================================================
    
    # 应用字体格式
    text_format_brush(element_id=element_id, font_size=font_size, font_name=font_name, 
                      color=color, bold=bold, italic=italic, underline=underline)
    
    # 【XML 修复】强制开启 "Shape to Fit Text"
    force_shape_to_fit_text(shape)
    
    # 【锁定宽度】防止 PPT 渲染时自动改变宽度
    shape.width = Inches(current_width_inch)

    #print(f"Set text content for element {element_id}. Lines: {total_lines}, New Height: {new_height_inch:.2f} in")
    return f"Set text content for element {element_id}"

import re
import ast


def append_text(element_id: str, text: str):
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    if not isinstance(text, str):
        text = str(text) if text is not None else ""

    text = text.replace('\\n', '\n').replace('\r\n', '\n').replace('\r', '\n')

    tf = shape.text_frame
    lines = text.split('\n')

    # --- 原有逻辑：样式源获取 ---
    style_source = None
    if tf.paragraphs:
        if not tf.paragraphs[-1].text.strip():
            last_p = tf.paragraphs[-1]
            # 这是一个空段落且没有 runs，移除它以避免空行
            if len(last_p.runs) == 0:
                p_element = last_p._p
                p_element.getparent().remove(p_element)

        if tf.paragraphs:
            style_source = tf.paragraphs[-1]

    if lines and lines[0] == '':
        lines = lines[1:]

    for line in lines:
        p = tf.add_paragraph()
        
        # 按照已有样式
        if style_source:
            p.level = style_source.level
            p.alignment = style_source.alignment

            ref_font = style_source.font
            if ref_font.size is None and style_source.runs:
                ref_font = style_source.runs[0].font

            if ref_font.size is not None:
                p.font.size = ref_font.size
            if ref_font.name is not None:
                p.font.name = ref_font.name
            if hasattr(ref_font, 'color') and hasattr(ref_font.color, 'rgb'):
                try:
                    p.font.color.rgb = ref_font.color.rgb
                except AttributeError:
                    pass
            p.font.bold = ref_font.bold
            p.font.italic = ref_font.italic

        # --- 内容处理 ---
        stripped_line = line.lstrip()

        # 根据内容覆盖层级 (Level)
        if stripped_line.startswith('•'):
            p.level = 0
        elif stripped_line.startswith('◦'):
            p.level = 1

        # 应用 Markdown 粗体并填入文本
        _apply_markdown_bold(p, line)


    calc_font_size_pt = 44.0
    if tf.paragraphs:
        last_p = tf.paragraphs[-1]
        if last_p.font.size:
            calc_font_size_pt = last_p.font.size.pt
        elif last_p.runs and last_p.runs[0].font.size:
            calc_font_size_pt = last_p.runs[0].font.size.pt
        elif style_source and style_source.font.size:
             calc_font_size_pt = style_source.font.size.pt

    #提取文本框内的【所有】文本 (旧 + 新)
    # 我们需要重构完整的文本字符串来计算总行数
    full_text_list = []
    for p in tf.paragraphs:
        full_text_list.append(p.text)
    full_text_content = "\n".join(full_text_list)

    # 3. 获取当前宽度 (英寸)
    current_width_inch = shape.width.inches
    
    # 4. 执行高度估算逻辑 (与 insert_textbox 保持一致)
    avg_char_width_pt = calc_font_size_pt * 0.55  
    
    # 计算有效宽度 (减去一点内边距 buffer)
    safe_width_pt = (current_width_inch * 72) - 10 
    if safe_width_pt < 10: safe_width_pt = 10
    
    chars_per_line = safe_width_pt / avg_char_width_pt
    if chars_per_line < 1: chars_per_line = 1

    # 计算总折行数
    total_lines = 0
    full_paragraphs = full_text_content.split('\n')
    
    for paragraph_text in full_paragraphs:
        length = len(paragraph_text)
        if length == 0:
            total_lines += 1
        else:
            total_lines += math.ceil(length / chars_per_line)

    line_height_pt = calc_font_size_pt * 1.2
    new_height_inch = ((total_lines * line_height_pt) / 72.0) + 0.2

    shape.height = Inches(new_height_inch)

    force_shape_to_fit_text(shape)
    shape.width = Inches(current_width_inch)

    return f"Appended text to element {element_id}."

def add_bullet_point(element_id: str, level: int = 0):
    shape = get_current_state().get_shape(element_id)

    if not hasattr(shape, 'text_frame'):
        return f"Error: Element {element_id} is not a text element"

    tf = shape.text_frame

    count = 0
    for p in tf.paragraphs:
        if not p.text.strip() and not p.runs:
            continue

        current_level = p.level if p.level is not None else level

        bullet = ""
        if current_level == 0:
            bullet = "• "
        elif current_level == 1:
            bullet = "◦ "

        current_text = p.text.lstrip()
        if current_text.startswith('•') or current_text.startswith('◦'):
            continue

        if p.runs:
            first_run = p.runs[0]
            first_run.text = f"{bullet}{first_run.text}"
        else:
            new_run = p.add_run()
            new_run.text = bullet + p.text

        count += 1

    return f"Added bullets to {count} paragraphs in element {element_id}"

from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml  # <--- 改用这个

def force_shape_to_fit_text(shape):
    """
    直接操作 XML，强制开启 'Shape to Fit Text' (高度自适应)
    并移除所有 'No Autofit' 或 'Text to Fit Shape' 的标记
    """
    # 获取文本框属性节点 (bodyPr)
    txBody = shape.text_frame._element
    bodyPr = txBody.find(qn('a:bodyPr'))

    # 1. 移除可能存在的冲突标签 (保持不变)
    # 移除 <a:noAutofit/> (固定大小)
    noAutofit = bodyPr.find(qn('a:noAutofit'))
    if noAutofit is not None:
        bodyPr.remove(noAutofit)
    
    # 移除 <a:normAutofit/> (缩小文字适应形状)
    normAutofit = bodyPr.find(qn('a:normAutofit'))
    if normAutofit is not None:
        bodyPr.remove(normAutofit)

    # 2. 强制插入 <a:spAutoFit/> (形状适应文字)
    # 检查是否已经存在，不存在则添加
    if bodyPr.find(qn('a:spAutoFit')) is None:
        # --- 核心修改点 ---
        # 使用 parse_xml 创建元素，并自动注入 'a' 命名空间
        spAutoFit = parse_xml('<a:spAutoFit %s/>' % nsdecls('a'))
        bodyPr.append(spAutoFit)

    # 3. 确保 wrap="square" (自动换行)
    # 如果 XML 里 wrap="none"，这里修正为 "square"
    if bodyPr.get('wrap') != 'square':
        bodyPr.set('wrap', 'square')



import math
def insert_textbox(left: float, top: float, width: float, height: float, text: str = "",
                   font_size: float = 44, font_name: str = "Arial", color: str = 'black',
                   bold: bool = False, italic: bool = False, underline: bool = False,
                   element_id: str = None):
    """
    插入文本框：宽度固定，高度根据【纯英文】内容估算并自动设置
    """
    if not isinstance(text, str):
        text = str(text) if text is not None else ""

    clean_text = text.replace('\\n', '\n')
    clean_text = clean_text.replace('\r\n', '\n') 
    clean_text = clean_text.replace('\r', '\n')
    
    avg_char_width_pt = font_size * 0.55
    textbox_width_pt = width * 72
    # 减去一点左右内边距 (0.1英寸 ≈ 7.2pt)
    safe_width_pt = textbox_width_pt - 7.2
    
    # 一行大约能放多少个字母
    chars_per_line = safe_width_pt / avg_char_width_pt
    if chars_per_line < 1: chars_per_line = 1

    # 2. 计算总行数
    total_lines = 0
    # 先按强制换行符分割
    paragraphs = text.split('\n')
    
    for p in paragraphs:
        length = len(p)
        if length == 0:
            total_lines += 1 # 空行也占位
        else:
            total_lines += math.ceil(length / chars_per_line)


    line_height_pt = font_size * 1.2
    estimated_height_pt = total_lines * line_height_pt
    
    # 转换回英寸，并加上一点上下边距 buffer (0.2英寸)
    final_height_inch = (estimated_height_pt / 72.0) + 0.2

    final_height_inch = max(height, final_height_inch)

    # 使用计算出的 final_height_inch 创建文本框
    textbox = get_current_state().slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(final_height_inch)
    )
    
    textbox.height = Inches(final_height_inch)

    tf = textbox.text_frame

    # 清理并填入文本
    tf.clear()
    tf.word_wrap = True 
    
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if not p.text.strip() and len(p.runs) == 0:
            p._p.getparent().remove(p._p)

    lines = text.split('\n')
    for line in lines:
        p = tf.add_paragraph()
        stripped_line = line.lstrip()
        if stripped_line.startswith('•'):
            p.level = 0
        elif stripped_line.startswith('◦'):
            p.level = 1
        _apply_markdown_bold(p, line)

    # 注册 ID
    if not element_id:
        element_id = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1
    get_current_state().shape_map[element_id] = textbox
    textbox.name = element_id

    # 应用格式
    text_format_brush(element_id=element_id, font_size=font_size, font_name=font_name, 
                      color=color, bold=bold, italic=italic, underline=underline)

    # 垂直对齐
    tf.vertical_anchor = MSO_ANCHOR.TOP

    force_shape_to_fit_text(textbox)
    
  
    textbox.width = Inches(width)

    print(f"Calculated Height: {textbox.height / 914400:.2f} inches (Lines: {total_lines})")
    
    return element_id

def delete_element(element_id: str):
    """
    删除元素

    Args:
        element_id: 元素ID
    """
    shape = get_current_state().get_shape(element_id)

    # 从 slide 中删除
    sp = shape.element
    sp.getparent().remove(sp)

    # 从映射表中删除
    del get_current_state().shape_map[element_id]

    return f"Deleted element {element_id}"


# ============================================================================
# 6. 图片操作 API
# ============================================================================

def insert_image(image_path: str, left: float, top: float,
                 width: Optional[float] = None, height: Optional[float] = None, element_id: Optional[str] = None):
    """
    插入图片

    Args:
        image_path: 图片文件路径
        left, top: 位置（厘米）
        width, height: 大小（厘米），None表示使用原始大小
    """
    if not os.path.isabs(image_path):
        # print(f"Image path {image_path} is not absolute, joining with pptx folder path: {get_current_state().pptx_folder_path}")
        image_path = os.path.join(get_current_state().pptx_folder_path, 'images_and_tables/'+image_path) # TODO

    picture = get_current_state().slide.shapes.add_picture(
        image_path,
        Inches(left), Inches(top),
        width=Inches(width) if width else None,
        height=Inches(height) if height else None
    )

    if not element_id:
        element_id = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1
    get_current_state().shape_map[element_id] = picture
    picture.name = element_id

    # print(f"Inserted image at {image_path} as element {element_id}")
    return f"{element_id}"


def replace_image(element_id: str, new_image_path: Optional[str] = None, new_image_element_id: Optional[str] = None):
    """
    替换图片（保持位置和大小）

    Args:
        element_id: 原图片元素ID
        new_image_path: 新图片路径
    """
    #TODO: to be tested
    if new_image_element_id:
        old_shape = get_current_state().get_shape(element_id)

        # 保存原始属性
        left = old_shape.left
        top = old_shape.top
        width = old_shape.width
        height = old_shape.height
        
        delete_element(element_id)
        set_element_position(new_image_element_id, left, top)
        set_element_size(new_image_element_id, width, height)   
        return f"Replaced image for element {element_id} with element {new_image_element_id}"
        
    if not os.path.isabs(new_image_path):
        new_image_path = os.path.join(get_current_state().pptx_folder_path, 'images_and_tables/'+new_image_path)
    old_shape = get_current_state().get_shape(element_id)

    # 保存原始属性
    left = old_shape.left
    top = old_shape.top
    width = old_shape.width
    height = old_shape.height

    # 删除旧图片
    delete_element(element_id)

    # 插入新图片
    picture = get_current_state().slide.shapes.add_picture(
        new_image_path, left, top, width, height
    )

    picture.name = element_id
    get_current_state().shape_map[element_id] = picture

    return f"Replaced image for element {element_id}"


# ============================================================================
# 工具函数
# ============================================================================

def get_element_info(element_id: str) -> Dict:
    """
    获取元素详细信息

    Args:
        element_id: 元素ID

    Returns:
        包含元素信息的字典
    """
    shape = get_current_state().get_shape(element_id)

    info = {
        "id": element_id,
        "name": shape.name,
        "type": str(shape.shape_type),
        "left": shape.left / 914400,
        "top": shape.top / 914400,
        "width": shape.width / 914400,
        "height": shape.height / 914400,
    }

    if hasattr(shape, 'text_frame'):
        info["text"] = shape.text_frame.text
        info["is_text"] = True
    else:
        info["is_text"] = False

    return info


def get_element_bounds(element_id: str) -> Tuple[float, float, float, float]:
    """
    获取元素边界

    Returns:
        (left, top, right, bottom) 以INCH为单位
    """
    shape = get_current_state().get_shape(element_id)

    left_inch = shape.left / 914400
    top_inch = shape.top / 914400
    width_inch = shape.width / 914400
    height_inch = shape.height / 914400

    return (left_inch, top_inch, left_inch + width_inch, top_inch + height_inch)


# ============================================================================
# 图例和标注 API
# ============================================================================


def add_callout(target_element_id: str, text: str,
                position: str = "right", offset: float = 1.5, element_id: Optional[str] = None):
    """
    为元素添加标注

    Args:
        target_element_id: 目标元素ID
        text: 标注文本
        position: 'left', 'right', 'top', 'bottom'
        offset: 距离目标的偏移量

    Returns:
        标注元素ID
    """
    target_shape = get_current_state().get_shape(target_element_id)

    # 计算标注位置
    target_left = target_shape.left / 914400
    target_top = target_shape.top / 914400
    target_width = target_shape.width / 914400
    target_height = target_shape.height / 914400

    callout_width = 4
    callout_height = 1.5

    if position == "right":
        callout_left = target_left + target_width + offset
        callout_top = target_top + target_height / 2 - callout_height / 2
    elif position == "left":
        callout_left = target_left - callout_width - offset
        callout_top = target_top + target_height / 2 - callout_height / 2
    elif position == "top":
        callout_left = target_left + target_width / 2 - callout_width / 2
        callout_top = target_top - callout_height - offset
    elif position == "bottom":
        callout_left = target_left + target_width / 2 - callout_width / 2
        callout_top = target_top + target_height + offset
    else:
        callout_left = target_left + target_width + offset
        callout_top = target_top

    if not element_id:
        element_id = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1

    # 创建标注文本框
    insert_textbox(callout_left, callout_top, callout_width, callout_height, text, font_size=44, element_id=element_id)
    set_shape_style(element_id, fill_color="light_yellow")
    set_line_style(element_id, "black", width=1)


# ============================================================================
# 批量操作 API
# ============================================================================

def batch_set_font_size(element_ids: List[str], size: float):
    """批量设置字体大小"""
    results = []
    for element_id in element_ids:
        result = set_text_font_size(element_id, size)
        results.append(result)
    return f"Batch set font size for {len(element_ids)} elements"


def batch_set_color(element_ids: List[str], color: str):
    """批量设置颜色"""
    results = []
    for element_id in element_ids:
        result = set_text_color(element_id, color)
        results.append(result)
    return f"Batch set color for {len(element_ids)} elements"


def batch_delete_elements(element_ids: List[str]):
    """批量删除元素"""
    for element_id in element_ids:
        delete_element(element_id)
    return f"Batch deleted {len(element_ids)} elements"


def move_group(element_ids: List[str], dx: Optional[float] = 0, dy: Optional[float] = 0):
    """Moves a list of elements by dx, dy."""
    for eid in element_ids:
        try:
            move_element_relative(eid, dx, dy)
        except Exception as e:
            print(f"Error moving element {eid}: {e}")

# TODO: element id
def clone_element(source_id: str, new_left: float, new_top: float,new_id: Optional[str] = None) -> str:
    """
    克隆元素到新位置

    Args:
        source_id: 源元素ID
        new_left, new_top: 新位置（厘米）
        new_name: 新元素名称

    Returns:
        新元素ID
    """
    source_shape = get_current_state().get_shape(source_id)

    # 获取源元素的XML
    source_element = source_shape.element

    # 复制元素
    from copy import deepcopy
    new_element = deepcopy(source_element)

    # 添加到slide
    source_shape.element.getparent().append(new_element)

    # 获取新shape对象
    new_shape = get_current_state().slide.shapes[-1]
    #print(new_shape)
    # 设置新位置
    new_shape.left = Inches(new_left)
    new_shape.top = Inches(new_top)

    if new_id:
        new_shape.name = new_id
    else:
        new_shape.name = str(get_current_state()._next_element_id)
        get_current_state()._next_element_id += 1
    element_id = new_shape.name
    get_current_state().shape_map[element_id] = new_shape
    return f"{element_id}"


def send_to_back_by_id(element_id: str):
    """
    Send shape backward by moving its XML element near the beginning of spTree,
    but AFTER required header nodes (p:nvGrpSpPr, p:grpSpPr).
    """
    shape = get_current_state().get_shape(element_id)
    spTree = shape.element.getparent()
    el = shape.element

    # Remove first (must remove before re-inserting)
    spTree.remove(el)

    # p:spTree usually starts with: p:nvGrpSpPr, p:grpSpPr
    # Insert after them if present; otherwise fallback to 0.
    insert_idx = 0
    try:
        children = list(spTree)
        if len(children) >= 2:
            if children[0].tag == qn("p:nvGrpSpPr") and children[1].tag == qn("p:grpSpPr"):
                insert_idx = 2
    except Exception:
        insert_idx = 0

    spTree.insert(insert_idx, el)

# ============================================================================
# 执行器 - 与 LLM 集成
# ============================================================================

def execute_api_calls(api_calls: List[str], prs, logger) -> str:
    """
    执行API调用列表

    Args:
        api_calls: API调用字符串列表
        prs: Presentation 对象

    Returns:
        错误信息（如果有）
    """
    # 从 prs 更新全局状态
    get_current_state().set_from_prs(prs)

    error_info = ""
    success_count = 0

    # 构建API上下文
    api_context = globals().copy()

    # 确保 shape_map 在上下文中可访问
    api_context['shape_map'] = get_current_state().shape_map
    api_context['slide'] = get_current_state().slide
    api_context['prs'] = prs

    logger.info("=" * 60)
    logger.info(f"Executing {len(api_calls)} API call(s)")
    logger.info(f"Shape map has {len(get_current_state().shape_map)} elements: {list(get_current_state().shape_map.keys())[:5]}...")
    logger.info("=" * 60)

    for i, line in enumerate(api_calls, 1):
        line = line.strip()
        if not line or line.startswith('#'):
            continue

        try:
            logger.info(f"\n[{i}] Executing: {line}")
            fixed_line = line.replace("\n", "\\n")
            result = eval(fixed_line, {"__builtins__": {}}, api_context)
            logger.info(f"    ✓ {result}")
            success_count += 1

        except Exception as e:
            error_msg = f"[Line {i}] {line}\n         Error: {type(e).__name__}: {e}"
            logger.info(f"    ✗ {error_msg}")
            error_info += f"{error_msg}\n\n"

    logger.info("\n" + "=" * 60)
    logger.info(
        f"Results: {success_count}/{len([l for l in api_calls if l.strip() and not l.strip().startswith('#')])} succeeded")
    logger.info("=" * 60)

    return error_info


# 兼容旧的接口名称
def API_executor(lines, api_context=None, prs=None, logger=None) -> str:
    """
    执行API调用（兼容接口）

    Args:
        lines: API调用字符串列表
        api_context: 忽略（为了兼容性保留）
        prs: Presentation 对象

    Returns:
        错误信息
    """
    if prs is None:
        raise ValueError("Must provide prs (Presentation object)")

    return execute_api_calls(lines, prs, logger)


if __name__ == "__main__":
    dic = {}
    from ..tools.pptx_parser import parse_pptx_to_json, parse_pptx_to_json_for_review
    # get_current_state() = PosterState()
    prs = get_current_state().prs = Presentation("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen copy.pptx")
    slide = get_current_state().slide = prs.slides[0]  # only process the first slide
    parse_pptx_to_json(pptx_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen copy.pptx",)
    
    print(get_current_state().shape_map.keys())
    print(get_current_state().prs.slides[0].shapes.element)
    count = 2
    for shape in slide.shapes:
        shape.name = str(shape.shape_id)  # Assign element_id to shape.name for easier tracking
        dic[shape.name] = shape
        # _element_to_shape_map[shape.name] = shape
        print(shape.name)
        
    # get_current_state() = PosterState()
    # get_current_state().load_poster(
    #     "/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen copy.pptx")
    # set_globalget_current_state()(get_current_state())

    # # from .image_tools import convert_pptx_to_png
    # # #set_text_content(element_id='32', text='Can Watermarked LLMs Be Identified by Users via Crafted Prompts?\\nA. Liu (Tsinghua University), S. Guan (Beijing University of Posts and Telecommunications), Y. Liu (Tsinghua University), L. Pan (Tsinghua University), Y. Zhang (The Chinese University of Hong Kong), L. Fang (University of Illinois at Chicago), L. Wen (Tsinghua University), P.S. Yu (University of Illinois at Chicago), X. Hu (Hongkong University of Science and Technology (Guangzhou))', font_size=100, font_name='Arial', color='purple', bold=False, italic=False, underline=False)
    # # insert_textbox(left=1.0, top=5.0, width=8.0, height=2.0, text='This poster is generated by PosterGen, an AI-powered poster generation system.\\nFor more information, visit: https://postergen.example.com', font_size=44, font_name='Arial', color='red', bold=False, italic=False, underline=False, element_id='100')
    insert_shape(left=0.5, top=4.5, width=9.0, height=2.5,shape_type='rectangle', fill_color='#FFD700', line_color='#FFD700', line_width=1.0, element_id='101')
    # # #move_element_relative('100', delta_x=0, delta_y=5)
    # # clone_element(source_id='32', new_left=1.0, new_top=8.0, new_id='102')
    # # # delete_element('101')
    # get_current_state().save_poster("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.pptx")
    
    prs = get_current_state().prs
    get_current_state().prs.save("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.pptx")
    
    
    # parse_pptx_to_json_for_review(pptx_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.pptx")
    from .image_tools import convert_pptx_to_png

    # # convert_pptx_to_png("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen copy 2.pptx", output_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.png")
    
    
    convert_pptx_to_png("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.pptx", output_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.png")
    # convert_pptx_to_png("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen.pptx", output_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen.png")
    # from .image_tools import convert_pptx_to_png
               
    # # convert_pptx_to_png("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Can_Watermarked_LLMs_be_Identified_by_Users_via_Crafted_Prompts/ByPosterGen_test2.pptx", output_path="/root/Poster-Edit/benchmark-v1-202511062200/Can_Watermarked_LLMs_be_Identified_by_Users_via_Crafted_Prompts/test_1.png")


    prs = get_current_state().prs = Presentation("/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/test1.pptx")
    # slide = get_current_state().slide = prs.slides[0]  # only process the first slide
    # # parse_pptx_to_json(pptx_path="/root/Poster_ppt_edit/Poster-Edit/benchmark-v1-202511062200/Revisiting_Robustness_in_Graph_Machine_Learning/ByPosterGen copy.pptx",)
    
    # print(get_current_state().shape_map.keys())
    # print(get_current_state().prs.slides[0].shapes.element)
    
    for shape in prs.slides[0].shapes:
        # shape.name = str(shape.shape_id)  # Assign element_id to shape.name for easier tracking
        # _element_to_shape_map[shape.name] = shape
        print(shape.name)
