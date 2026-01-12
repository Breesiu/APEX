import uuid
import json
import os
from typing import Dict
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.oxml import parse_xml
from ..schema import PosterJSON, PosterElement, TextRun

from typing import Dict, List, Set, Any
from ..schema import PosterJSON, PosterElement
from copy import deepcopy
from . import pptx_execuator

class PosterFilter: # TODO: need to compute excute time
    """Utility class for filtering poster JSON and APIs based on task requirements"""
    
    @staticmethod
    def _remove_empty_values(obj: Any) -> Any:
        """
        Recursively remove keys with empty values (None, [], "", {}) from dict/list.
        
        Args:
            obj: Dictionary, list, or other value
        
        Returns:
            Cleaned object with empty values removed
        """


        from pydantic import BaseModel
        
        # Handle Pydantic BaseModel objects
        if isinstance(obj, BaseModel):
            # Convert to dict first, then clean
            obj_dict = obj.model_dump() if hasattr(obj, 'model_dump') else obj.dict()
            # print(obj_dict)
            return PosterFilter._remove_empty_values(obj_dict)
        
        
        if isinstance(obj, dict):
            cleaned = {}
            for k, v in obj.items():
                if k == 'z_index' or k == 'image_path': #or k == 'runs':
                    continue # Skip z_index field
                if k == 'elements':
                    # print(v)
                    pass
                cleaned_v = PosterFilter._remove_empty_values(v)
                # Keep the value only if it's not empty
                if cleaned_v is not None and cleaned_v != [] and cleaned_v != "" and cleaned_v != {}:
                    cleaned[k] = cleaned_v
            # print(cleaned)
            return cleaned if cleaned else None
        elif isinstance(obj, list):
            cleaned_list = []
            for item in obj:
                # print(item)
                cleaned_item = PosterFilter._remove_empty_values(item)
                if cleaned_item is not None and cleaned_item != [] and cleaned_item != "" and cleaned_item != {}:
                    cleaned_list.append(cleaned_item)
            # print(cleaned_list)
            return cleaned_list if cleaned_list else None
        else:
            # print(obj)
            # Return the value as-is (including 0, False, etc.)
            return obj
        

    @staticmethod
    def filter_revised_json(original_json: PosterJSON, revised_json: PosterJSON, all_revised_field=False) -> PosterJSON:
        filtered = {
            "slide_width": revised_json.slide_width,
            "slide_height": revised_json.slide_height,
            "elements": []
            }
        filtered_fields = ['id', 'type', 'left', 'top', 'width', 'height', 'text', 'fill_color', 'border_color', 'border_width', 'border_shape']
        for element in revised_json.elements:
            has_changes = []
                        
            # Find corresponding element in original JSON
            orig_elem = next((e for e in original_json.elements if e.id == element.id), None)
            if orig_elem is None: 
                # has_changes = filtered_fields.copy()
                has_changes = ['left', 'top', 'width', 'height']
            # Compare fields to see if any changes
            for field in filtered_fields:
                if getattr(element, field, None) != getattr(orig_elem, field, None):
                    if field in ['left', 'top', 'width', 'height']:
                        has_changes.extend(['left', 'top', 'width', 'height']) # only consider position changes as a whole
                    if all_revised_field:
                        has_changes.append(field)
            
            if has_changes:
                has_changes.extend(['id', 'type'])
                dic = {key: getattr(element, key, None) for key in has_changes}
                element_filter = PosterElement(**dic)
                filtered['elements'].append(element_filter)
        return filtered
    
# # Global mapping: element_id -> shape_id for tracking
# _element_to_shape_map: Dict[str, int] = {}

def _emu_to_inch(emu_value: int) -> float: # TODO
    """Convert EMU (English Metric Units) to inches"""
    return float(emu_value) / 914400

def _rgb_to_hex(rgb) -> str: # TODO
    """Convert RGB tuple to hex color"""
    if rgb is None:
        return None
    try:
        return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except:
        return None

def _get_theme_colors(prs) -> Dict[int, str]:
    """Extract theme color mapping from presentation"""
    theme_colors = {}
    try:
        slide_master = prs.slide_masters[0]
        
        theme_part = None
        theme_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        
        for rel in slide_master.part.rels.values():
             if rel.reltype == theme_rel_type:
                 theme_part = rel.target_part
                 break
        
        if theme_part:
            theme_element = parse_xml(theme_part.blob)
            namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            clrScheme = theme_element.find('.//a:clrScheme', namespaces)
            
            if clrScheme is not None:
                tag_mapping = {
                    'dk1': MSO_THEME_COLOR.DARK_1,
                    'lt1': MSO_THEME_COLOR.LIGHT_1,
                    'dk2': MSO_THEME_COLOR.DARK_2,
                    'lt2': MSO_THEME_COLOR.LIGHT_2,
                    'accent1': MSO_THEME_COLOR.ACCENT_1,
                    'accent2': MSO_THEME_COLOR.ACCENT_2,
                    'accent3': MSO_THEME_COLOR.ACCENT_3,
                    'accent4': MSO_THEME_COLOR.ACCENT_4,
                    'accent5': MSO_THEME_COLOR.ACCENT_5,
                    'accent6': MSO_THEME_COLOR.ACCENT_6,
                    'hlink': MSO_THEME_COLOR.HYPERLINK,
                    'folHlink': MSO_THEME_COLOR.FOLLOWED_HYPERLINK
                }
                
                for child in clrScheme:
                    tag = child.tag.split('}')[-1]
                    if tag in tag_mapping:
                        color_hex = None
                        srgbClr = child.find('.//a:srgbClr', namespaces)
                        sysClr = child.find('.//a:sysClr', namespaces)
                        
                        if srgbClr is not None:
                            color_hex = srgbClr.get('val')
                        elif sysClr is not None:
                            color_hex = sysClr.get('lastClr')
                        
                        if color_hex:
                            theme_colors[tag_mapping[tag]] = f"#{color_hex}"
    except Exception:
        pass
        
    return theme_colors

def _get_effective_font_style(run, paragraph):
    """
    Resolve font style by checking run properties, then paragraph default properties.
    Returns a dict with 'size', 'bold', 'italic', 'underline', 'color_hex', 'name'.
    """
    style = {
        'size': None,
        'bold': None,
        'italic': None,
        'underline': None,
        'color_hex': None,
        'name': None
    }
    
    # 1. Check Run Properties (Direct Formatting)
    if run.font.size:
        style['size'] = run.font.size.pt
    
    if run.font.bold is not None:
        style['bold'] = run.font.bold
        
    if run.font.italic is not None:
        style['italic'] = run.font.italic
        
    if run.font.underline is not None:
        style['underline'] = run.font.underline
        
    if run.font.name is not None:
        style['name'] = run.font.name
        
    if hasattr(run.font, 'color') and run.font.color:
        try:
            style['color_hex'] = _rgb_to_hex(run.font.color.rgb)
        except:
            pass

    # 2. Check Paragraph Default Properties (defRPr) if anything is missing
    # Accessing internal lxml element: paragraph._p.pPr.defRPr
    try:
        pPr = paragraph._p.pPr
        if pPr is not None:
            defRPr = pPr.defRPr
            if defRPr is not None:
                # Font Size
                if style['size'] is None and defRPr.sz is not None:
                    # sz is in centipoints (1/100 pt)
                    style['size'] = float(defRPr.sz) / 100.0
                
                # Bold
                if style['bold'] is None and defRPr.b is not None:
                    style['bold'] = defRPr.b
                
                # Italic
                if style['italic'] is None and defRPr.i is not None:
                    style['italic'] = defRPr.i
                    
                # Underline
                if style['underline'] is None and defRPr.u is not None:
                     style['underline'] = (defRPr.u != 'none')

                # Font Name (latin typeface)
                if style['name'] is None:
                    # Check for latin typeface
                    if hasattr(defRPr, 'latin'):
                        style['name'] = defRPr.latin.typeface

                # Color
                if style['color_hex'] is None:
                    # defRPr can have solidFill/srgbClr
                    # We need to traverse children of defRPr to find solidFill -> srgbClr
                    if hasattr(defRPr, 'solidFill'):
                         if defRPr.solidFill.srgbClr is not None:
                             val = defRPr.solidFill.srgbClr.val
                             if val:
                                 style['color_hex'] = f"#{val}"
    except Exception:
        pass
        
    return style

def parse_pptx_to_json(pptx_path: str) -> PosterJSON:
    """
    Parse PPTX file and extract all element properties as JSON.
    
    Args:
        pptx_path: Path to PPTX file
    
    Returns:
        PosterJSON object with all elements
    """
    # global _element_to_shape_map
    # _element_to_shape_map.clear()
    _element_to_shape_map = {}
    prs = pptx_execuator.get_current_state().prs = Presentation(pptx_path)
    slide = pptx_execuator.get_current_state().slide = prs.slides[0]  # only process the first slide
    
    # Extract theme colors
    theme_color_map = _get_theme_colors(prs)
    
    elements = []
    
    element_counter = 2

    for shape in pptx_execuator.get_current_state().slide.shapes:
        shape.name = str(shape.shape_id)  # Assign element_id to shape.name for easier tracking   #TODO: for multi-step editing, reorder all elements
        _element_to_shape_map[shape.name] = shape
        #print(shape.name)
        # cNvPr = shape._element.get_or_add_cNvPr()
        # cNvPr.set("descr", str(shape.shape_id))
        # shape.element.set('descr', str(shape.shape_id))
        # shape._element.set('mytag', str(shape.shape_id))
        # Determine element type
        elem_type = "shape"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            elem_type = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            elem_type = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            elem_type = "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            elem_type = "line"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                elem_type = "rectangle"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                elem_type = "textbox"
        elif (hasattr(shape, 'has_text_frame') and shape.has_text_frame) or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elem_type = "textbox"
        
        # Base properties
        elem = PosterElement(
            id=shape.name,
            # id=shape.name.split('_')[-1],
            # id=shape._element.get_or_add_cNvPr().get("descr"),
            # id=shape.element.get('descr'),
            # id=shape._element.get('mytag'),
            # id=str(shape.shape_id),
            type=elem_type,
            left=round(_emu_to_inch(shape.left), 2),
            top=round(_emu_to_inch(shape.top), 2),
            width=round(_emu_to_inch(shape.width), 2),
            height=round(_emu_to_inch(shape.height), 2)
        )
        #print(shape.name)
        # Extract text content
        if (elem_type == "textbox" or elem_type == "rectangle") and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            elem.text = shape.text_frame.text
            
            for para in shape.text_frame.paragraphs:
                bullet_level = para.level if hasattr(para, 'level') else None
                
                for run in para.runs:
                    # Use helper to resolve inherited styles
                    style = _get_effective_font_style(run, para)
                    if elem.main_font_size is None:
                        elem.main_font_size = style['size'] if style['size'] is not None else elem.main_font_size
                    text_run = TextRun(
                        text=run.text,
                        bold=style['bold'],
                        italic=style['italic'],
                        underline=style['underline'],
                        font_name=style['name'],
                        font_size=style['size'],
                        font_color=style['color_hex'],
                        bullet_level=bullet_level
                    )
                    elem.runs.append(text_run)
        
        # Extract shape fill
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:  # SOLID
                    try:
                        elem.fill_color = _rgb_to_hex(shape.fill.fore_color.rgb)
                    except AttributeError:
                        if hasattr(shape.fill.fore_color, 'type') and shape.fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                             # Try to resolve theme color
                             theme_color_idx = shape.fill.fore_color.theme_color
                             if theme_color_idx in theme_color_map:
                                 elem.fill_color = theme_color_map[theme_color_idx]
                             else:
                                 elem.fill_color = str(theme_color_idx)
            except Exception:
                pass
        
        # Extract border info
        if hasattr(shape, 'line'):
            try:
                # Check if line is solid (1) to ensure it's visible
                if shape.line.fill.type == 1:
                    # Extract width first (independent of color)
                    if shape.line.width:
                        elem.border_width = round(_emu_to_inch(shape.line.width), 2)
                    
                    # Try to extract color (might fail for Theme colors)
                    try:
                        elem.border_color = _rgb_to_hex(shape.line.color.rgb)
                    except AttributeError:
                         if hasattr(shape.line.color, 'type') and shape.line.color.type == MSO_COLOR_TYPE.SCHEME:
                             # Try to resolve theme color
                             theme_color_idx = shape.line.color.theme_color
                             if theme_color_idx in theme_color_map:
                                 elem.border_color = theme_color_map[theme_color_idx]
                             else:
                                 elem.border_color = str(theme_color_idx)
            except Exception:
                pass

        # Extract image path
        if elem_type == "picture":
            try:
                image = shape.image
                elem.image_path = image.filename if hasattr(image, 'filename') else None
            except:
                pass
        
        elements.append(elem)
        
    pptx_execuator.get_current_state().shape_map = _element_to_shape_map

    # # Reassign IDs starting from 1
    # # Create a temporary copy of the mapping to look up shape_ids by old_id
    # temp_map = _element_to_shape_map.copy()
    # _element_to_shape_map.clear() # Clear old mapping
    
    # id_mapping = {}
    
    # for new_id, elem in enumerate(elements, start=2):
    #     str_id = str(new_id)
    #     old_id = elem.id
    #     id_mapping[old_id] = str_id
        
    #     elem.id = str_id
        
    #     # Rebuild the map: element_id -> shape_id
    #     if old_id in temp_map:
    #         _element_to_shape_map[str_id] = temp_map[old_id]

    poster_json = PosterJSON(
        slide_width=round(_emu_to_inch(prs.slide_width),2),
        slide_height=round(_emu_to_inch(prs.slide_height),2),
        elements=elements
    )
    
    # Save debug info
    try:
        base_dir = os.path.dirname(pptx_path)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        
    #     mapping_path = os.path.join(base_dir, f"{base_name}_id_mapping.json")
        json_path = os.path.join(base_dir, f"{base_name}_poster_layout.json")
        
    #     # with open(mapping_path, 'w', encoding='utf-8') as f:
    #     #     json.dump(id_mapping, f, indent=2)
            
        # with open(json_path, 'w', encoding='utf-8') as f:
        #     try:
        #         f.write(PosterFilter._remove_empty_values(poster_json).model_dump_json(indent=2))
        #         print(PosterFilter._remove_empty_values(poster_json))
        #     except Exception as e:
        #         print(PosterFilter._remove_empty_values(poster_json))
        #         f.write(PosterJSON(**PosterFilter._remove_empty_values(poster_json)).model_dump_json(exclude_unset=True))#, exclude={'elements': {'__all__': {'runs': True}}}))
        #         # f.write(PosterJSON(**PosterFilter._remove_empty_values(poster_json)).model_dump_json(indent=2))
    except Exception as e:
        print(f"Warning: Failed to save debug files: {e}")
                
    # except Exception as e:
    #     print(f"Warning: Failed to save debug files: {e}")
    return PosterJSON(**PosterFilter._remove_empty_values(poster_json))
    # return poster_json
#indent=2,


def parse_pptx_to_json_for_review(pptx_path: str) -> PosterJSON:
    """
    Parse PPTX file and extract all element properties as JSON.
    
    Args:
        pptx_path: Path to PPTX file
    
    Returns:
        PosterJSON object with all elements
    """
    # global _element_to_shape_map
    # _element_to_shape_map.clear()
    _element_to_shape_map = {}
    prs = pptx_execuator.get_current_state().prs = Presentation(pptx_path)
    slide = pptx_execuator.get_current_state().slide = prs.slides[0]  # only process the first slide
    
    # Extract theme colors
    theme_color_map = _get_theme_colors(prs)
    
    elements = []
    
    element_counter = 2

    for shape in slide.shapes:
        _element_to_shape_map[shape.name] = shape # NOTE: difference with parse_pptx_to_json
        # Assign element_id to shape.name for easier tracking
        # print(shape.name)
        # Determine element type
        elem_type = "shape"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            elem_type = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            elem_type = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            elem_type = "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            elem_type = "line"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                elem_type = "rectangle"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                elem_type = "textbox"
        elif (hasattr(shape, 'has_text_frame') and shape.has_text_frame) or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elem_type = "textbox"
        
        # Base properties
        elem = PosterElement(
            id=shape.name,
            type=elem_type,
            left=round(_emu_to_inch(shape.left), 2),
            top=round(_emu_to_inch(shape.top), 2),
            width=round(_emu_to_inch(shape.width), 2),
            height=round(_emu_to_inch(shape.height), 2)
        )
        #print(shape.name)
        # Extract text content
        if (elem_type == "textbox" or elem_type == "rectangle") and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            elem.text = shape.text_frame.text
            
            for para in shape.text_frame.paragraphs:
                bullet_level = para.level if hasattr(para, 'level') else None
                
                for run in para.runs:
                    # Use helper to resolve inherited styles
                    style = _get_effective_font_style(run, para)
                    if elem.main_font_size is None:
                        elem.main_font_size = style['size'] if style['size'] is not None else elem.main_font_size
                    text_run = TextRun(
                        text=run.text,
                        bold=style['bold'],
                        italic=style['italic'],
                        underline=style['underline'],
                        font_name=style['name'],
                        font_size=style['size'],
                        font_color=style['color_hex'],
                        bullet_level=bullet_level
                    )
                    elem.runs.append(text_run)
        
        # Extract shape fill
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:  # SOLID
                    try:
                        elem.fill_color = _rgb_to_hex(shape.fill.fore_color.rgb)
                    except AttributeError:
                        if hasattr(shape.fill.fore_color, 'type') and shape.fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                             # Try to resolve theme color
                             theme_color_idx = shape.fill.fore_color.theme_color
                             if theme_color_idx in theme_color_map:
                                 elem.fill_color = theme_color_map[theme_color_idx]
                             else:
                                 elem.fill_color = str(theme_color_idx)
            except:
                pass
        
        # Extract border info
        if hasattr(shape, 'line'):
            try:
                # Check if line is solid (1) to ensure it's visible
                if shape.line.fill.type == 1:
                    # Extract width first (independent of color)
                    if shape.line.width:
                        elem.border_width = round(_emu_to_inch(shape.line.width), 2)
                    
                    # Try to extract color (might fail for Theme colors)
                    try:
                        elem.border_color = _rgb_to_hex(shape.line.color.rgb)
                    except AttributeError:
                         if hasattr(shape.line.color, 'type') and shape.line.color.type == MSO_COLOR_TYPE.SCHEME:
                             # Try to resolve theme color
                             theme_color_idx = shape.line.color.theme_color
                             if theme_color_idx in theme_color_map:
                                 elem.border_color = theme_color_map[theme_color_idx]
                             else:
                                 elem.border_color = str(theme_color_idx)
            except:
                pass
        
        # Extract image path
        if elem_type == "picture":
            try:
                image = shape.image
                elem.image_path = image.filename if hasattr(image, 'filename') else None
            except:
                pass
        
        elements.append(elem)
        
    pptx_execuator.get_current_state().shape_map = _element_to_shape_map
    
    
    # prs.save(pptx_path)
    # # Reassign IDs starting from 1
    # # Create a temporary copy of the mapping to look up shape_ids by old_id
    # temp_map = _element_to_shape_map.copy()
    # _element_to_shape_map.clear() # Clear old mapping
    
    # id_mapping = {}
    
    # for new_id, elem in enumerate(elements, start=2):
    #     str_id = str(new_id)
    #     old_id = elem.id
    #     id_mapping[old_id] = str_id
        
    #     elem.id = str_id
        
    #     # Rebuild the map: element_id -> shape_id
    #     if old_id in temp_map:
    #         _element_to_shape_map[str_id] = temp_map[old_id]

    poster_json = PosterJSON(
        slide_width=round(_emu_to_inch(prs.slide_width),2),
        slide_height=round(_emu_to_inch(prs.slide_height),2),
        elements=elements
    )
    
    # Save debug info
    try:
        base_dir = os.path.dirname(pptx_path)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        
    #     mapping_path = os.path.join(base_dir, f"{base_name}_id_mapping.json")
        json_path = os.path.join(base_dir, f"{base_name}_poster_layout.json")
        
    #     # with open(mapping_path, 'w', encoding='utf-8') as f:
    #     #     json.dump(id_mapping, f, indent=2)
            
        with open(json_path, 'w', encoding='utf-8') as f:
            try:
                f.write(PosterFilter._remove_empty_values(poster_json).model_dump_json(indent=2))
                print(PosterFilter._remove_empty_values(poster_json))
            except Exception as e:
                print(PosterFilter._remove_empty_values(poster_json))
                f.write(PosterJSON(**PosterFilter._remove_empty_values(poster_json)).model_dump_json(exclude_unset=True))#, exclude={'elements': {'__all__': {'runs': True}}}))
                # f.write(PosterJSON(**PosterFilter._remove_empty_values(poster_json)).model_dump_json(indent=2))
    except Exception as e:
        print(f"Warning: Failed to save debug files: {e}")
                
    # except Exception as e:
    #     print(f"Warning: Failed to save debug files: {e}")
    return PosterJSON(**PosterFilter._remove_empty_values(poster_json))
    # return poster_json
#indent=2,



    

if __name__ == "__main__":
    poster_json = parse_pptx_to_json_for_review("/root/Poster-Edit/benchmark-v1-202511062200/Revisiting Robustness in Graph Machine Learning/ByPosterGen.pptx")
    # print(poster_json)


def get_element_shape_map() -> Dict[str, int]:
    """Return the current element ID to shape ID mapping"""
    return 
    # return _element_to_shape_map.copy()

